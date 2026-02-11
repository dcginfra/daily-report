"""Unit tests for report_data, format_markdown, and format_slides.

Tests dataclass construction, Markdown output, PPTX slide deck generation,
and CLI flag validation for the slides export feature.

Run with: python3 -m pytest tests/test_formatters.py -v
"""

import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation

from daily_report.report_data import (
    AuthoredPR,
    ReportData,
    ReviewedPR,
    SummaryStats,
    WaitingPR,
)
from daily_report.format_markdown import format_markdown
from daily_report.format_slides import format_slides

PROJECT_ROOT = Path(__file__).resolve().parent.parent


# ---------------------------------------------------------------------------
# Helpers: reusable fixtures / factories
# ---------------------------------------------------------------------------

def _make_report(**kwargs) -> ReportData:
    """Create a ReportData with sensible defaults, overridden by kwargs."""
    defaults = dict(
        user="testuser",
        date_from="2026-02-10",
        date_to="2026-02-10",
        authored_prs=[],
        reviewed_prs=[],
        waiting_prs=[],
        summary=SummaryStats(
            total_prs=0, repo_count=0, merged_count=0,
            open_count=0, themes=[], is_range=False,
        ),
    )
    defaults.update(kwargs)
    return ReportData(**defaults)


def _make_full_report() -> ReportData:
    """Report with all sections populated across two repos."""
    return _make_report(
        user="alice",
        date_from="2026-02-10",
        date_to="2026-02-10",
        authored_prs=[
            AuthoredPR(
                repo="org/alpha", title="Add login", number=10,
                status="Open", additions=50, deletions=10,
                contributed=False, original_author=None,
            ),
            AuthoredPR(
                repo="org/beta", title="Fix crash", number=20,
                status="Merged", additions=0, deletions=0,
                contributed=True, original_author="bob",
            ),
        ],
        reviewed_prs=[
            ReviewedPR(
                repo="org/alpha", title="Update docs", number=11,
                author="charlie", status="Open",
            ),
        ],
        waiting_prs=[
            WaitingPR(
                repo="org/beta", title="Refactor DB", number=21,
                reviewers=["dave", "eve"], created_at="2026-02-08",
                days_waiting=2,
            ),
        ],
        summary=SummaryStats(
            total_prs=4, repo_count=2, merged_count=1,
            open_count=2, themes=["feat", "fix"], is_range=False,
        ),
    )


# ---------------------------------------------------------------------------
# ReportData smoke tests
# ---------------------------------------------------------------------------

class TestReportDataConstruction:
    """Verify dataclass construction with known values."""

    def test_authored_pr_fields(self):
        pr = AuthoredPR(
            repo="org/repo", title="Add feature", number=1,
            status="Open", additions=10, deletions=5,
            contributed=False, original_author=None,
        )
        assert pr.repo == "org/repo"
        assert pr.number == 1
        assert pr.status == "Open"
        assert pr.additions == 10
        assert pr.contributed is False

    def test_reviewed_pr_fields(self):
        pr = ReviewedPR(
            repo="org/repo", title="Fix bug", number=2,
            author="user1", status="Merged",
        )
        assert pr.author == "user1"
        assert pr.status == "Merged"

    def test_waiting_pr_fields(self):
        pr = WaitingPR(
            repo="org/repo", title="New API", number=3,
            reviewers=["r1", "r2"], created_at="2026-02-01",
            days_waiting=5,
        )
        assert pr.reviewers == ["r1", "r2"]
        assert pr.days_waiting == 5

    def test_summary_stats_fields(self):
        s = SummaryStats(
            total_prs=10, repo_count=3, merged_count=4,
            open_count=6, themes=["feat", "fix"], is_range=True,
        )
        assert s.total_prs == 10
        assert s.is_range is True

    def test_report_data_defaults(self):
        r = ReportData(user="u", date_from="2026-01-01", date_to="2026-01-01")
        assert r.authored_prs == []
        assert r.reviewed_prs == []
        assert r.waiting_prs == []
        assert r.summary.total_prs == 0

    def test_report_data_full(self):
        r = _make_full_report()
        assert r.user == "alice"
        assert len(r.authored_prs) == 2
        assert len(r.reviewed_prs) == 1
        assert len(r.waiting_prs) == 1
        assert r.summary.total_prs == 4


# ---------------------------------------------------------------------------
# format_markdown tests
# ---------------------------------------------------------------------------

class TestFormatMarkdownSingleDate:
    """Single date report with authored, reviewed, and waiting PRs."""

    def test_header_single_date(self):
        report = _make_full_report()
        md = format_markdown(report)
        assert "# Daily Report" in md
        assert "2026-02-10" in md
        assert ".." not in md

    def test_authored_pr_present(self):
        report = _make_full_report()
        md = format_markdown(report)
        assert "Add login #10" in md
        assert "`org/alpha`" in md

    def test_reviewed_pr_present(self):
        report = _make_full_report()
        md = format_markdown(report)
        assert "Update docs #11" in md
        assert "(charlie)" in md

    def test_waiting_pr_present(self):
        report = _make_full_report()
        md = format_markdown(report)
        assert "Refactor DB #21" in md
        assert "**dave**" in md
        assert "**eve**" in md
        assert "2 days" in md

    def test_summary_merged_today(self):
        report = _make_full_report()
        md = format_markdown(report)
        assert "merged today" in md
        assert "4 PRs across 2 repos" in md

    def test_themes_in_summary(self):
        report = _make_full_report()
        md = format_markdown(report)
        assert "feat, fix" in md


class TestFormatMarkdownDateRange:
    """Date range report (is_range=True)."""

    def test_header_range(self):
        report = _make_report(
            date_from="2026-02-03", date_to="2026-02-09",
            summary=SummaryStats(
                total_prs=5, repo_count=2, merged_count=3,
                open_count=2, themes=[], is_range=True,
            ),
        )
        md = format_markdown(report)
        assert "2026-02-03 .. 2026-02-09" in md

    def test_summary_merged_not_today(self):
        report = _make_report(
            date_from="2026-02-03", date_to="2026-02-09",
            summary=SummaryStats(
                total_prs=5, repo_count=2, merged_count=3,
                open_count=2, themes=[], is_range=True,
            ),
        )
        md = format_markdown(report)
        assert "merged today" not in md
        assert "3 merged" in md


class TestFormatMarkdownEmpty:
    """Empty report with no PRs."""

    def test_no_authored_message(self):
        report = _make_report()
        md = format_markdown(report)
        assert "No authored or contributed PRs" in md

    def test_no_reviewed_message(self):
        report = _make_report()
        md = format_markdown(report)
        assert "No reviewed or approved PRs" in md

    def test_no_waiting_message(self):
        report = _make_report()
        md = format_markdown(report)
        assert "No PRs waiting for review" in md

    def test_themes_default(self):
        report = _make_report()
        md = format_markdown(report)
        assert "general development" in md


class TestFormatMarkdownContributed:
    """Contributed PRs show original author."""

    def test_contributed_shows_author(self):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="Contrib PR", number=99,
                    status="Merged", additions=0, deletions=0,
                    contributed=True, original_author="original_dev",
                ),
            ],
        )
        md = format_markdown(report)
        assert "(original_dev)" in md

    def test_non_contributed_no_author(self):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="My PR", number=88,
                    status="Open", additions=5, deletions=2,
                    contributed=False, original_author=None,
                ),
            ],
        )
        md = format_markdown(report)
        # No parenthetical author name before the status
        assert "My PR #88" in md
        # There should be no extra parens around the PR title/number
        lines = [l for l in md.splitlines() if "#88" in l]
        assert len(lines) == 1
        # No author parens between title and status
        assert "()" not in lines[0]


class TestFormatMarkdownAdditionsDeletions:
    """Open/Draft PRs show +additions/-deletions; Merged/Closed do not."""

    def test_open_shows_stats(self):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="Open PR", number=1,
                    status="Open", additions=30, deletions=12,
                    contributed=False, original_author=None,
                ),
            ],
        )
        md = format_markdown(report)
        assert "+30" in md
        assert "12" in md

    def test_draft_shows_stats(self):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="Draft PR", number=2,
                    status="Draft", additions=15, deletions=3,
                    contributed=False, original_author=None,
                ),
            ],
        )
        md = format_markdown(report)
        assert "+15" in md

    def test_merged_no_stats(self):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="Merged PR", number=3,
                    status="Merged", additions=0, deletions=0,
                    contributed=False, original_author=None,
                ),
            ],
        )
        md = format_markdown(report)
        line = [l for l in md.splitlines() if "#3" in l][0]
        assert "+" not in line or "+0" not in line
        # Merged should not have the stats suffix
        assert "(+0" not in line

    def test_closed_no_stats(self):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="Closed PR", number=4,
                    status="Closed", additions=0, deletions=0,
                    contributed=False, original_author=None,
                ),
            ],
        )
        md = format_markdown(report)
        line = [l for l in md.splitlines() if "#4" in l][0]
        assert "(+0" not in line


# ---------------------------------------------------------------------------
# format_slides tests
# ---------------------------------------------------------------------------

class TestFormatSlidesBasic:
    """Basic slide deck generation with full report."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        self.report = _make_full_report()
        self.output_path = str(tmp_path / "report.pptx")
        format_slides(self.report, self.output_path)
        self.prs = Presentation(self.output_path)

    def test_file_created(self):
        assert Path(self.output_path).exists()

    def test_slide_count(self):
        # 1 title + 2 project slides (org/alpha, org/beta) + 1 summary = 4
        assert len(self.prs.slides) == 4

    def test_title_slide_title(self):
        title_slide = self.prs.slides[0]
        assert title_slide.shapes.title.text == "Activity Report"

    def test_title_slide_contains_user(self):
        title_slide = self.prs.slides[0]
        subtitle = title_slide.placeholders[1].text
        assert "alice" in subtitle

    def test_title_slide_contains_date(self):
        title_slide = self.prs.slides[0]
        subtitle = title_slide.placeholders[1].text
        assert "2026-02-10" in subtitle

    def test_project_slides_repo_names(self):
        # Projects sorted: org/alpha, org/beta
        repo_titles = [
            self.prs.slides[i].shapes.title.text for i in range(1, 3)
        ]
        assert repo_titles == ["org/alpha", "org/beta"]

    def test_summary_slide_title(self):
        summary_slide = self.prs.slides[-1]
        assert summary_slide.shapes.title.text == "Summary"

    def test_summary_slide_contains_metrics(self):
        summary_slide = self.prs.slides[-1]
        text = summary_slide.placeholders[1].text
        assert "Total PRs: 4" in text
        assert "Repositories: 2" in text
        assert "1 merged today" in text
        assert "2 still open" in text
        assert "feat, fix" in text


class TestFormatSlidesEmpty:
    """Empty report produces title + summary only (no project slides)."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        self.report = _make_report()
        self.output_path = str(tmp_path / "empty.pptx")
        format_slides(self.report, self.output_path)
        self.prs = Presentation(self.output_path)

    def test_slide_count(self):
        # 1 title + 0 projects + 1 summary = 2
        assert len(self.prs.slides) == 2

    def test_title_slide(self):
        assert self.prs.slides[0].shapes.title.text == "Activity Report"

    def test_summary_slide(self):
        assert self.prs.slides[-1].shapes.title.text == "Summary"


class TestFormatSlidesSorted:
    """Projects are sorted alphabetically."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/zebra", title="Z PR", number=1,
                    status="Open", additions=1, deletions=0,
                    contributed=False, original_author=None,
                ),
                AuthoredPR(
                    repo="org/alpha", title="A PR", number=2,
                    status="Open", additions=1, deletions=0,
                    contributed=False, original_author=None,
                ),
                AuthoredPR(
                    repo="org/middle", title="M PR", number=3,
                    status="Open", additions=1, deletions=0,
                    contributed=False, original_author=None,
                ),
            ],
            summary=SummaryStats(
                total_prs=3, repo_count=3, merged_count=0,
                open_count=3, themes=[], is_range=False,
            ),
        )
        self.output_path = str(tmp_path / "sorted.pptx")
        format_slides(report, self.output_path)
        self.prs = Presentation(self.output_path)

    def test_project_order(self):
        # Skip title (0), take project slides (1,2,3), skip summary (4)
        repo_titles = [
            self.prs.slides[i].shapes.title.text for i in range(1, 4)
        ]
        assert repo_titles == ["org/alpha", "org/middle", "org/zebra"]


class TestFormatSlidesDateRange:
    """Title slide shows date range when is_range is used."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        report = _make_report(
            user="rangeuser",
            date_from="2026-02-03",
            date_to="2026-02-09",
            summary=SummaryStats(
                total_prs=0, repo_count=0, merged_count=0,
                open_count=0, themes=[], is_range=True,
            ),
        )
        self.output_path = str(tmp_path / "range.pptx")
        format_slides(report, self.output_path)
        self.prs = Presentation(self.output_path)

    def test_title_slide_range(self):
        subtitle = self.prs.slides[0].placeholders[1].text
        assert "2026-02-03 .. 2026-02-09" in subtitle

    def test_summary_merged_not_today(self):
        text = self.prs.slides[-1].placeholders[1].text
        assert "merged today" not in text
        assert "merged" in text


class TestFormatSlidesSectionOmission:
    """Empty sections are omitted from project slides."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        # Only authored PRs, no reviewed or waiting
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="Solo PR", number=1,
                    status="Open", additions=5, deletions=2,
                    contributed=False, original_author=None,
                ),
            ],
            summary=SummaryStats(
                total_prs=1, repo_count=1, merged_count=0,
                open_count=1, themes=[], is_range=False,
            ),
        )
        self.output_path = str(tmp_path / "section.pptx")
        format_slides(report, self.output_path)
        self.prs = Presentation(self.output_path)

    def test_has_authored_section(self):
        project_slide = self.prs.slides[1]
        text = project_slide.placeholders[1].text
        assert "Authored" in text

    def test_no_reviewed_section(self):
        project_slide = self.prs.slides[1]
        text = project_slide.placeholders[1].text
        assert "Reviewed" not in text

    def test_no_waiting_section(self):
        project_slide = self.prs.slides[1]
        text = project_slide.placeholders[1].text
        assert "Waiting" not in text


class TestFormatSlidesContributedPR:
    """Contributed PRs show original author in slide text."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="Contrib", number=99,
                    status="Merged", additions=0, deletions=0,
                    contributed=True, original_author="origdev",
                ),
            ],
            summary=SummaryStats(
                total_prs=1, repo_count=1, merged_count=1,
                open_count=0, themes=[], is_range=False,
            ),
        )
        self.output_path = str(tmp_path / "contrib.pptx")
        format_slides(report, self.output_path)
        self.prs = Presentation(self.output_path)

    def test_original_author_in_slide(self):
        project_slide = self.prs.slides[1]
        text = project_slide.placeholders[1].text
        assert "(origdev)" in text


class TestFormatSlidesOpenVsMergedStats:
    """Open/Draft show +/-; Merged/Closed do not."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        report = _make_report(
            authored_prs=[
                AuthoredPR(
                    repo="org/repo", title="OpenPR", number=1,
                    status="Open", additions=30, deletions=12,
                    contributed=False, original_author=None,
                ),
                AuthoredPR(
                    repo="org/repo", title="MergedPR", number=2,
                    status="Merged", additions=0, deletions=0,
                    contributed=False, original_author=None,
                ),
            ],
            summary=SummaryStats(
                total_prs=2, repo_count=1, merged_count=1,
                open_count=1, themes=[], is_range=False,
            ),
        )
        self.output_path = str(tmp_path / "stats.pptx")
        format_slides(report, self.output_path)
        self.prs = Presentation(self.output_path)

    def test_open_pr_has_stats(self):
        project_slide = self.prs.slides[1]
        paragraphs = project_slide.placeholders[1].text_frame.paragraphs
        open_line = [p.text for p in paragraphs if "OpenPR" in p.text][0]
        assert "+30/-12" in open_line

    def test_merged_pr_no_stats(self):
        project_slide = self.prs.slides[1]
        paragraphs = project_slide.placeholders[1].text_frame.paragraphs
        merged_line = [p.text for p in paragraphs if "MergedPR" in p.text][0]
        assert "+0" not in merged_line
        assert "/-" not in merged_line


# ---------------------------------------------------------------------------
# CLI flag tests
# ---------------------------------------------------------------------------

class TestCLISlidesFlags:
    """CLI argument validation for --slides and --slides-output."""

    def test_slides_output_without_slides_errors(self):
        """--slides-output without --slides should error."""
        cmd = [
            sys.executable, "-m", "daily_report",
            "--slides-output", "/tmp/test.pptx",
        ]
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=30,
            cwd=str(PROJECT_ROOT),
        )
        assert result.returncode != 0
        assert "--slides-output requires --slides" in result.stderr
