"""PPTX slide deck formatter for daily-report.

Requires python-pptx: pip install python-pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from daily_report.report_data import (
    ReportData, AuthoredPR, ReviewedPR, WaitingPR,
)


def format_slides(report: ReportData, output_path: str) -> None:
    """Render the report as a PPTX slide deck.

    Args:
        report: Complete report data.
        output_path: File path to write the .pptx file.

    Raises:
        OSError: If the file cannot be written (permissions, missing directory).
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, report)

    projects = _group_by_repo(report)
    for repo_name in sorted(projects):
        group = projects[repo_name]
        _add_project_slide(
            prs, repo_name,
            group["authored"], group["reviewed"], group["waiting"],
        )

    _add_summary_slide(prs, report)

    prs.save(output_path)


# --- internal helpers (private) ---


def _add_title_slide(prs: Presentation, report: ReportData) -> None:
    """Add the title slide with user and date range."""
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Activity Report"
    if report.date_from == report.date_to:
        subtitle_text = f"{report.user}\n{report.date_from}"
    else:
        subtitle_text = f"{report.user}\n{report.date_from} .. {report.date_to}"
    slide.placeholders[1].text = subtitle_text


def _add_project_slide(prs: Presentation, repo_name: str,
                        authored: list[AuthoredPR],
                        reviewed: list[ReviewedPR],
                        waiting: list[WaitingPR]) -> None:
    """Add a project slide with grouped bullet lists."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = repo_name

    tf = slide.placeholders[1].text_frame
    tf.clear()
    first_paragraph = True

    if authored:
        p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
        first_paragraph = False
        p.text = "Authored / Contributed"
        p.level = 0
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(14)

        for pr in authored:
            p = tf.add_paragraph()
            p.text = _authored_pr_text(pr)
            p.level = 1
            p.runs[0].font.size = Pt(12)

    if reviewed:
        p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
        first_paragraph = False
        p.text = "Reviewed"
        p.level = 0
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(14)

        for pr in reviewed:
            p = tf.add_paragraph()
            p.text = _reviewed_pr_text(pr)
            p.level = 1
            p.runs[0].font.size = Pt(12)

    if waiting:
        p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
        first_paragraph = False
        p.text = "Waiting for Review"
        p.level = 0
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(14)

        for pr in waiting:
            p = tf.add_paragraph()
            p.text = _waiting_pr_text(pr)
            p.level = 1
            p.runs[0].font.size = Pt(12)


def _add_summary_slide(prs: Presentation, report: ReportData) -> None:
    """Add the summary slide with aggregate metrics."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Summary"

    s = report.summary
    themes_str = ", ".join(s.themes) if s.themes else "general development"
    merged_label = "merged" if s.is_range else "merged today"

    bullets = [
        f"Total PRs: {s.total_prs}",
        f"Repositories: {s.repo_count}",
        f"{s.merged_count} {merged_label}",
        f"{s.open_count} still open",
        f"Key themes: {themes_str}",
    ]

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.runs[0].font.size = Pt(14)


def _group_by_repo(report: ReportData) -> dict[str, dict]:
    """Group all PR lists by repository name.

    Returns:
        Dict mapping repo name to {"authored": [...], "reviewed": [...], "waiting": [...]}.
        Only repos with at least one item are included.
    """
    projects: dict[str, dict] = {}
    for pr in report.authored_prs:
        projects.setdefault(pr.repo, {"authored": [], "reviewed": [], "waiting": []})
        projects[pr.repo]["authored"].append(pr)
    for pr in report.reviewed_prs:
        projects.setdefault(pr.repo, {"authored": [], "reviewed": [], "waiting": []})
        projects[pr.repo]["reviewed"].append(pr)
    for pr in report.waiting_prs:
        projects.setdefault(pr.repo, {"authored": [], "reviewed": [], "waiting": []})
        projects[pr.repo]["waiting"].append(pr)
    return projects


def _authored_pr_text(pr: AuthoredPR) -> str:
    """Build bullet text for an authored/contributed PR."""
    text = f"{pr.title} #{pr.number}"
    if pr.contributed and pr.original_author:
        text += f" ({pr.original_author})"
    text += f" -- {pr.status}"
    if pr.status in ("Open", "Draft"):
        text += f" (+{pr.additions}/-{pr.deletions})"
    return text


def _reviewed_pr_text(pr: ReviewedPR) -> str:
    """Build bullet text for a reviewed PR."""
    return f"{pr.title} #{pr.number} ({pr.author}) -- {pr.status}"


def _waiting_pr_text(pr: WaitingPR) -> str:
    """Build bullet text for a PR waiting for review."""
    reviewers = ", ".join(pr.reviewers)
    return f"{pr.title} #{pr.number} -- reviewer: {reviewers} -- {pr.days_waiting} days"
